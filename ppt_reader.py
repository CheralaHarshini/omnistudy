from pptx import Presentation

def read_ppt(file):

    prs=Presentation(file)

    text=""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text+=shape.text

    return text