import streamlit as st
import pdfplumber
from pptx import Presentation
from transformers import pipeline

st.title("OmniStudy AI")

st.write("Upload a PDF or PPT to get summary and important questions.")

@st.cache_resource
def load_model():
    model = pipeline(
        "text2text-generation",
        model="google/flan-t5-small"
    )
    return model

model = load_model()


def read_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t
    return text


def read_ppt(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def summarize(text):
    text = text[:2000]
    prompt = f"Summarize the following study material:\n{text}"
    result = model(prompt, max_length=150)
    return result[0]["generated_text"]


def generate_questions(text):
    prompt = f"Generate 5 important exam questions from this topic:\n{text}"
    result = model(prompt, max_length=150)
    return result[0]["generated_text"]


uploaded_file = st.file_uploader(
    "Upload PDF or PPT",
    type=["pdf", "pptx"]
)

if uploaded_file:

    if uploaded_file.name.endswith(".pdf"):
        text = read_pdf(uploaded_file)
    else:
        text = read_ppt(uploaded_file)

    st.subheader("Summary")

    summary = summarize(text)

    st.write(summary)

    st.subheader("Important Questions")

    questions = generate_questions(summary)

    st.write(questions)